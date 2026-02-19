from pptx import Presentation
from pptx.util import Inches, Pt
from text_rules.sentence_splitter import split_sentences
import sys
sys.path.append(".")

# ---------------- PATHS ----------------
want_bg = True

sahil_bg_path = "/Users/mayankjain/Documents/Sahil Jain/book_summary_slides/Books summary/Utility/Sahil jain logo page .png"
mayank_bg_path = "/Users/mayankjain/Desktop/Mayank_logo.png"
bg_path = sahil_bg_path

file_name = "change_your_life"
input_path = f"/Users/mayankjain/Documents/Sahil Jain/book_summary_slides/Books summary/raw/{file_name}.txt"
output_path = f"/Users/mayankjain/Documents/Sahil Jain/book_summary_slides/Books summary/curated/{file_name}.pptx"

# ---------------- CONFIG ----------------
MAX_WORDS_PER_SLIDE = 30
FONT_NAME = "Helvetica" 
FONT_SIZE = Pt(36) 
FONT_BOLD = True

LEFT_MARGIN = Inches(1.3)
WIDTH = Inches(10.3)
BOX_HEIGHT = Inches(1.4)

TOP_MARGIN = Inches(1.3)
BOTTOM_MARGIN = Inches(1.6)
# ----------------------------------------


# ---------------- HELPERS ----------------
def word_count(s):
    return len(s.split())

def is_heading(s):
    words = s.split()
    return (
        3 <= len(words) <= 10
        and not s.strip().endswith(('.', '?', '!'))
        and sum(w[0].isupper() for w in words if w[0].isalpha()) >= len(words) * 0.6
    )

def is_short_fragment(s):
    return word_count(s) <= 5

def is_long_sentence(s):
    return word_count(s) >= 22
# -----------------------------------------


# ---------------- READ TEXT ----------------
with open(input_path, "r", encoding="utf-8") as file:
    text = file.read()

sentences = split_sentences(text)


# ---------------- PACKING LOGIC ----------------
slides_chunks = []
i = 0

while i < len(sentences):
    s = sentences[i].strip()

    # HEADING
    if is_heading(s):
        slides_chunks.append([s])
        i += 1
        continue

    # LONG SENTENCE
    if is_long_sentence(s):
        slides_chunks.append([s])
        i += 1
        continue

    # SHORT FRAGMENTS (group up to 3)
    if is_short_fragment(s):
        group = [s]
        j = i + 1

        while j < len(sentences) and is_short_fragment(sentences[j]) and len(group) < 3:
            group.append(sentences[j])
            j += 1

        slides_chunks.append(group)
        i = j
        continue

    # NORMAL PACKING BY WORD LIMIT
    group = [s]
    total_words = word_count(s)
    j = i + 1

    while j < len(sentences):
        next_s = sentences[j]

        if total_words + word_count(next_s) > MAX_WORDS_PER_SLIDE:
            break

        group.append(next_s)
        total_words += word_count(next_s)
        j += 1

    slides_chunks.append(group)
    i = j


# ---------------- CREATE PPT ----------------
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_h = prs.slide_height

top_pos = TOP_MARGIN
middle_pos = (slide_h / 2) - (BOX_HEIGHT / 2)
bottom_pos = slide_h - BOX_HEIGHT - BOTTOM_MARGIN


for chunk in slides_chunks:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    if want_bg:
        slide.shapes.add_picture(
            bg_path, 0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    # Limit to max 3 lines visually
    if len(chunk) > 3:
        merged = " ".join(chunk[2:])
        chunk = [chunk[0], chunk[1], merged]

    # HEADING SLIDE
    if len(chunk) == 1 and is_heading(chunk[0]):
        box = slide.shapes.add_textbox(LEFT_MARGIN, middle_pos, WIDTH, BOX_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = chunk[0]
        p.font.size = Pt(46)
        p.font.bold = True
        p.font.name = FONT_NAME
        continue

    # POSITIONING
    if len(chunk) == 1:
        positions = [top_pos]
    elif len(chunk) == 2:
        positions = [top_pos, bottom_pos]
    else:
        positions = [top_pos, middle_pos, bottom_pos]

    # ADD TEXT
    for idx, line in enumerate(chunk):
        box = slide.shapes.add_textbox(LEFT_MARGIN, positions[idx], WIDTH, BOX_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = line
        p.font.size = FONT_SIZE
        p.font.bold = FONT_BOLD
        p.font.name = FONT_NAME


prs.save(output_path)
print("File saved into:", output_path)
