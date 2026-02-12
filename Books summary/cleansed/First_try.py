from pptx import Presentation
from pptx.util import Inches, Pt
from text_rules.sentence_splitter import split_sentences
import sys
sys.path.append(".")


bg_path = "/Users/mayankjain/Documents/Sahil Jain/Books summary/Utility/Sahil jain logo page .png"

file_name = "its_sound_crazy"
input_path = f"/Users/mayankjain/Documents/Sahil Jain/Books summary/raw/{file_name}.txt"
output_path = f"/Users/mayankjain/Documents/Sahil Jain/Books summary/curated/{file_name}.pptx"

with open(input_path, "r", encoding="utf-8") as file:
    text = file.read()

sentences = split_sentences(text)


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

left_margin = Inches(1.3)
width = Inches(10.3)

top_positions = [Inches(1.5), Inches(4.0)]

for i in range(0, len(sentences), 2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background image covering full slide
    slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    for idx in range(2):
        if i + idx < len(sentences):
            box = slide.shapes.add_textbox(left_margin, top_positions[idx], width, Inches(1.4))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sentences[i + idx]
            p.font.name = "Verdana"
            p.font.size = Pt(36)  # reduced font size to prevent collision

prs.save(output_path)
output_path
